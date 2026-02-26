import streamlit as st
import tempfile
import chromadb
from sentence_transformers import SentenceTransformer
from pypdf import PdfReader
from pptx import Presentation

st.set_page_config(page_title="Study Bot", layout="centered")
st.title("📚 Study Bot — No LLM / No API")

embedder = SentenceTransformer("all-MiniLM-L6-v2")

client = chromadb.PersistentClient(path="./db")
collection = client.get_or_create_collection("notes")

# -------- loaders --------

def load_pdf(path):
    reader = PdfReader(path)
    return "\n".join(p.extract_text() for p in reader.pages if p.extract_text())

def load_ppt(path):
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

# -------- upload --------

files = st.file_uploader(
    "Upload study files",
    type=["pdf", "pptx", "txt"],
    accept_multiple_files=True
)

if files:
    for f in files:
        with tempfile.NamedTemporaryFile(delete=False) as tmp:
            tmp.write(f.read())
            path = tmp.name

        if f.name.endswith(".pdf"):
            text = load_pdf(path)
        elif f.name.endswith(".pptx"):
            text = load_ppt(path)
        else:
            text = open(path).read()

        chunks = [text[i:i+700] for i in range(0, len(text), 700)]
        vecs = embedder.encode(chunks).tolist()

        collection.add(
            documents=chunks,
            embeddings=vecs,
            ids=[f"{f.name}_{i}" for i in range(len(chunks))]
        )

        st.success(f"{f.name} indexed")

# -------- query --------

q = st.text_input("Ask from your notes")

if q:
    qv = embedder.encode([q]).tolist()
    res = collection.query(query_embeddings=qv, n_results=3)

    st.subheader("📌 Best Matching Passages")
    for i, doc in enumerate(res["documents"][0], 1):
        st.markdown(f"**Match {i}:**")
        st.write(doc)
